#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(content_list):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(12)

def add_two_column_slide(title, left_content, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column
    txBox_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5))
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(8)

    # Right column
    txBox_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.5))
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(8)

# === SLIDES ===

# 1. Title
add_title_slide("TSN Switch Manager Dashboard", "KETI 전자부품연구원")

# 2. 목차
add_content_slide("목차", [
    "1. TSN 기술 개요",
    "2. 시스템 구성",
    "3. 대시보드 기능 소개",
    "4. 개발 결과물",
    "5. 시연"
])

# 3. TSN 개요
add_content_slide("TSN (Time-Sensitive Networking) 이란?", [
    "• IEEE 802.1 표준 기반의 결정론적 이더넷 기술",
    "• 기존 이더넷에 시간 동기화, 트래픽 스케줄링, 대역폭 보장 기능 추가",
    "• 산업 자동화, 자동차, 오디오/비디오 스트리밍 분야에 적용",
    "",
    "핵심 특징:",
    "  - 저지연: 마이크로초 단위의 예측 가능한 지연",
    "  - 시간 동기화: 네트워크 전체 1μs 이하 동기화",
    "  - 대역폭 보장: 중요 트래픽에 대한 대역폭 예약",
    "  - 공존성: 일반 이더넷 트래픽과 공존 가능"
])

# 4. TSN 핵심 프로토콜
add_content_slide("TSN 핵심 프로토콜", [
    "1. PTP (Precision Time Protocol) - IEEE 802.1AS",
    "   • 네트워크 전체 시간 동기화",
    "   • Grandmaster(GM) 클럭 기준으로 모든 노드 동기화",
    "   • 나노초(ns) 단위 정밀도",
    "",
    "2. TAS (Time-Aware Shaper) - IEEE 802.1Qbv",
    "   • 시간 기반 트래픽 스케줄링",
    "   • Gate Control List(GCL)로 전송 시점 제어",
    "",
    "3. CBS (Credit-Based Shaper) - IEEE 802.1Qav",
    "   • 크레딧 기반 대역폭 제어",
    "   • SR Class A/B 트래픽에 대역폭 보장"
])

# 5. PTP 동작 원리
add_two_column_slide("PTP 동작 원리", [
    "동기화 과정:",
    "1. GM이 Sync 메시지 전송",
    "2. Slave가 Delay Request 전송",
    "3. 경로 지연(Path Delay) 계산",
    "4. 클럭 오프셋 보정",
    "",
    "토폴로지:",
    "  Grandmaster ←→ Slave",
    "   (Board 1)      (Board 2)"
], [
    "주요 파라미터:",
    "",
    "• Clock Offset",
    "  클럭 차이 (목표: ±50ns 이내)",
    "",
    "• Path Delay",
    "  네트워크 전송 지연 (~5μs)",
    "",
    "• Sync Interval",
    "  동기화 주기 (125ms)"
])

# 6. TAS 동작 원리
add_content_slide("TAS 동작 원리", [
    "Gate Control List (GCL):",
    "• 8개 Traffic Class (TC0~TC7)",
    "• 각 슬롯에서 열림/닫힘 상태 제어",
    "• 1초 사이클, 8개 슬롯 × 125ms",
    "",
    "트래픽 클래스:",
    "  TC0: BE(BG) - Best Effort (Background)",
    "  TC5: Voice - 음성 트래픽 (SR Class A)",
    "  TC6: Video - 영상 트래픽 (SR Class B)",
    "  TC7: NC - Network Control"
])

# 7. CBS 동작 원리
add_two_column_slide("CBS 동작 원리", [
    "파라미터:",
    "",
    "• idleSlope",
    "  대기 시 크레딧 충전 속도",
    "",
    "• sendSlope",
    "  전송 시 크레딧 소비 속도",
    "",
    "• 크레딧 > 0 일 때만 전송 가능"
], [
    "대역폭 할당 예시:",
    "",
    "• TC5 (Voice): 250 Mbps (25%)",
    "",
    "• TC6 (Video): 375 Mbps (37.5%)",
    "",
    "• Best Effort: 375 Mbps (37.5%)"
])

# 8. 시스템 구성
add_two_column_slide("시스템 구성", [
    "하드웨어:",
    "",
    "• TSN 스위치 보드",
    "  - 칩셋: Microchip LAN9692",
    "  - 포트: 9포트 TSN 지원",
    "  - 속도: 1Gbps",
    "",
    "• 네트워크 토폴로지",
    "  Board #1 (GM) ↔ Board #2 (Slave)"
], [
    "소프트웨어:",
    "",
    "• 웹 대시보드",
    "  - Frontend: HTML5, CSS3, JS",
    "  - 배포: GitHub Pages",
    "  - 실시간 업데이트: 8Hz",
    "",
    "• 주요 기능",
    "  - PTP, TAS, CBS 모니터링",
    "  - 디바이스 상태 확인"
])

# 9. PTP Dashboard
add_content_slide("PTP Dashboard", [
    "실시간 모니터링 항목:",
    "• Clock Offset: GM과의 시간 차이 (ns)",
    "• Path Delay: 네트워크 전송 지연 (μs)",
    "• Steps Removed: GM으로부터의 홉 수",
    "• Clock Class: IEEE 1588 클럭 등급",
    "",
    "그래프:",
    "• Clock Offset History (실시간 그래프)",
    "• 200 샘플 유지, 8Hz 업데이트",
    "",
    "상태: Locked (|offset| < 50ns) / Acquiring"
])

# 10. TAS Dashboard
add_content_slide("TAS Dashboard", [
    "Gate Control List (GCL) 매트릭스:",
    "• 8개 슬롯 × 8개 Traffic Class",
    "• 각 셀: 게이트 열림(●) / 닫힘(○)",
    "",
    "GCL Timeline:",
    "• 1초 사이클 시각화",
    "• TC별 게이트 열림 구간 표시",
    "",
    "Traffic Test:",
    "• 모든 TC에 트래픽 발생",
    "• 지연시간, 지터 측정"
])

# 11. CBS Dashboard
add_content_slide("CBS Dashboard", [
    "Queue Configuration:",
    "• 8개 큐 (TC0~TC7) 설정 표시",
    "• idleSlope, sendSlope 값",
    "",
    "Credit Level Monitor:",
    "• TC5, TC6 크레딧 변화 그래프",
    "• 실시간 업데이트",
    "",
    "Bandwidth Test:",
    "• 지정 대역폭으로 트래픽 발생",
    "• 달성률, 손실률 측정"
])

# 12. Device Status
add_content_slide("Device Status", [
    "디바이스 카드:",
    "• 보드별 상태 (Online/Offline)",
    "• 역할, IP, MAC, 포트 수, Uptime",
    "",
    "Network Topology:",
    "• PTP 계층 구조 시각화",
    "• GM → Slave 연결 표시",
    "",
    "Port Status:",
    "• 포트별 링크 상태",
    "• RX/TX 바이트, 패킷 카운트"
])

# 13. 개발 결과물
add_two_column_slide("개발 결과물", [
    "완료 항목:",
    "✓ 웹 기반 TSN 대시보드 UI",
    "✓ PTP 동기화 시뮬레이션",
    "✓ TAS GCL 시각화",
    "✓ CBS 크레딧 모니터링",
    "✓ 디바이스 상태 모니터링",
    "✓ 실시간 그래프 (Canvas)",
    "✓ Retina 디스플레이 지원",
    "✓ GitHub Pages 배포"
], [
    "기술 스택:",
    "",
    "• Frontend",
    "  HTML5, CSS3, JavaScript",
    "",
    "• 그래프",
    "  Canvas API",
    "",
    "• 배포",
    "  GitHub Pages"
])

# 14. 접속 정보
add_content_slide("접속 정보", [
    "",
    "웹 대시보드 URL:",
    "https://hwkim3330.github.io/keti-tsn-cli-ui-sim/",
    "",
    "",
    "GitHub Repository:",
    "https://github.com/hwkim3330/keti-tsn-cli-ui-sim"
])

# 15. 시연
add_content_slide("시연 시나리오", [
    "1. PTP Dashboard",
    "   • 보드 선택 → 오프셋 값 확인",
    "   • 실시간 그래프 동작 확인",
    "",
    "2. TAS Dashboard",
    "   • GCL 매트릭스 확인",
    "   • Traffic Test 실행",
    "",
    "3. CBS Dashboard",
    "   • Credit Monitor 그래프 확인",
    "   • Bandwidth Test 실행",
    "",
    "4. Device Status",
    "   • 토폴로지 및 포트 상태 확인"
])

# 16. Thank you
add_title_slide("감사합니다", "Q&A")

# Save
prs.save('/home/jhchoi/keti-tsn-cli-ui-sim/TSN_Dashboard.pptx')
print("PPT created: TSN_Dashboard.pptx")
